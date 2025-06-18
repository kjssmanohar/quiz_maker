import streamlit as st
import requests
from PyPDF2 import PdfReader
from pptx import Presentation
import json
import re
import google.generativeai as genai

# === CONFIG ===
GEMINI_API_KEY = "AIzaSyBe9wzZboiPOSDMkp1XiycGmN-kFgDe814"
WEB_APP_URL = "https://script.google.com/macros/s/AKfycbxcmumXquoOFtuC0TRXBtbUHRk_xPbsxraYj7nsinHXzh1OyQXgvercY4Pjyk34U8tR/exec"  # 🔁 Replace with your deployed Apps Script URL

# === File Parsers ===
def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

# === Gemini AI Quiz Generator ===
def generate_quiz_questions(text, api_key):
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-2.0-flash")
    prompt = f"""
You are a quiz-making assistant.

From the content below, generate exactly 10 multiple-choice questions.
Each question must have:
- A "question" string
- An "options" list of 4 strings
- An "answer" string that exactly matches one of the options

Only output a valid JSON array in this format:
[
  {{
    "question": "Sample?",
    "options": ["A", "B", "C", "D"],
    "answer": "B"
  }},
  ...
]

Content:
{text[:2000]}
"""
    response = model.generate_content(prompt)
    return response.text.strip()

# === JSON Extractor ===
def extract_json(text):
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        matches = re.findall(r"\[\s*{.?}\s]", text, re.DOTALL)
        for match in matches:
            try:
                return json.loads(match)
            except:
                continue
        start = text.find('[')
        end = text.rfind(']')
        if start != -1 and end != -1:
            try:
                return json.loads(text[start:end+1])
            except:
                pass
        raise ValueError("❌ Could not extract valid JSON array.")

# === UI ===
st.title("📄 Auto Quiz Generator & Google Form Creator")

uploaded_file = st.file_uploader("Upload PDF or PPTX", type=["pdf", "pptx"])

if uploaded_file and st.button("Generate Quiz"):
    try:
        # 1. Extract text
        if uploaded_file.type == "application/pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_text_from_pptx(uploaded_file)
        else:
            st.error("Unsupported file type.")
            st.stop()

        # 2. Generate quiz via Gemini
        st.info("Generating quiz via Gemini...")
        quiz_raw = generate_quiz_questions(text, GEMINI_API_KEY)

        with st.expander("🔍 Gemini Output (Debug)"):
            st.code(quiz_raw, language="json")

        quiz_data = extract_json(quiz_raw)

        # 3. Send to Google Apps Script
        st.info("Creating Google Form...")
        payload = {
            "title": "Auto-Generated Quiz",
            "questions": quiz_data
        }

        response = requests.post(WEB_APP_URL, json=payload)
        if response.status_code != 200:
            st.error(f"❌ Form creation failed: {response.status_code} - {response.text}")
            st.stop()

        form_url = response.json().get("url")
        if not form_url:
            st.error(f"❌ No form URL returned: {response.json()}")
            st.stop()

        st.success(f"✅ Google Form ready! [Open Quiz Form]({form_url})")

    except Exception as e:
        st.error(f"❌ Error: {e}")

# === Top Scorer ===
st.markdown("---")
st.subheader("🏆 Top Scorer from Google Form Responses")

top_response = requests.get(f"{WEB_APP_URL}?mode=top_scorer")
try:
    top_data = top_response.json()
    if "error" in top_data:
        st.error(f"❌ {top_data['error']}")
    else:
        st.info(f"🥇 Top Scorer: *{top_data['name']}* — *{top_data['score']} points*")
except:
    st.error("❌ Failed to parse top scorer data.")